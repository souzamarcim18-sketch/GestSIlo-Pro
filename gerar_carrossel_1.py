"""
Gerador do Carrossel 1 (apresentação GestSilo) para Instagram.

Cria um PowerPoint editável, formato quadrado (1080x1080 a 100dpi),
seguindo a identidade visual do kit de marca GestSilo:
  - Cores primárias: #323E30 (verde escuro), #738D45 (verde campo)
  - Fonte: Inter (kit de marca indica "Inter SemiBold").
    O PowerPoint substitui automaticamente se a fonte nao estiver instalada.

Todo o texto fica em caixas editaveis. Sem imagens, icones ou emojis.
Slides 1 e 6 trazem um placeholder de foto opcional (retangulo tracejado).

Uso:
    python gerar_carrossel_1.py
"""

import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    print("python-pptx nao encontrado. Instale com: pip install python-pptx")
    sys.exit(1)


# ---------------------------------------------------------------------------
# Constantes de marca
# ---------------------------------------------------------------------------
VERDE_ESCURO = RGBColor(0x32, 0x3E, 0x30)   # #323E30
VERDE_CAMPO = RGBColor(0x73, 0x8D, 0x45)    # #738D45
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)         # #FFFFFF

FONTE = "Arial"  # fonte segura e editavel (kit indica "Inter SemiBold")

LADO = Inches(10.8)          # slide quadrado
MARGEM = Inches(0.9)         # respiro generoso nas bordas
LARGURA_CONTEUDO = LADO - (MARGEM * 2)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _definir_fundo(slide, cor):
    """Pinta o fundo do slide com uma cor solida."""
    fundo = slide.background
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = cor


def _set_run(run, texto, tamanho, cor, negrito=False):
    run.text = texto
    run.font.size = Pt(tamanho)
    run.font.bold = negrito
    run.font.name = FONTE
    run.font.color.rgb = cor


def adicionar_caixa(
    slide,
    left,
    top,
    width,
    height,
    paragrafos,
    cor,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    espacamento_paragrafos=Pt(12),
):
    """
    Adiciona uma caixa de texto editavel.

    `paragrafos` e uma lista de dicts: {"texto", "tamanho", "negrito"?, "cor"?}.
    Cada item vira um paragrafo.
    """
    caixa = slide.shapes.add_textbox(left, top, width, height)
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    for i, item in enumerate(paragrafos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i > 0:
            p.space_before = espacamento_paragrafos
        run = p.add_run()
        _set_run(
            run,
            item["texto"],
            item["tamanho"],
            item.get("cor", cor),
            item.get("negrito", False),
        )
    return caixa


def adicionar_logo(slide, cor_texto):
    """Caixa '[LOGO GESTSILO]' no canto inferior direito (tamanho 12)."""
    largura = Inches(2.6)
    altura = Inches(0.4)
    left = LADO - largura - Inches(0.4)
    top = LADO - altura - Inches(0.35)
    caixa = slide.shapes.add_textbox(left, top, largura, altura)
    tf = caixa.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    _set_run(run, "[LOGO GESTSILO]", 12, cor_texto, negrito=False)


def adicionar_placeholder_foto(slide, cor_texto):
    """
    Retangulo vazio (sem preenchimento, borda fina tracejada) na area
    central-inferior + caixa de texto '[ESPACO PARA FOTO OPCIONAL]'.
    """
    largura = Inches(7.0)
    altura = Inches(2.6)
    left = (LADO - largura) / 2
    top = Inches(7.0)

    forma = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, largura, altura
    )
    forma.fill.background()  # sem preenchimento

    linha = forma.line
    linha.color.rgb = cor_texto
    linha.width = Pt(1.0)
    # borda tracejada
    ln = linha._get_or_add_ln()
    dash = ln.find(qn("a:prstDash"))
    if dash is None:
        dash = ln.makeelement(qn("a:prstDash"), {})
        ln.append(dash)
    dash.set("val", "dash")

    # remove o texto padrao da autoshape e usa caixa sobreposta
    forma.text_frame.text = ""

    caixa = slide.shapes.add_textbox(left, top, largura, altura)
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, "[ESPACO PARA FOTO OPCIONAL]", 16, cor_texto, negrito=False)


# ---------------------------------------------------------------------------
# Montagem da apresentacao
# ---------------------------------------------------------------------------
def construir():
    prs = Presentation()
    prs.slide_width = LADO
    prs.slide_height = LADO
    layout_branco = prs.slide_layouts[6]  # layout em branco

    fundos = []  # registro para o relatorio final

    # ----- SLIDE 1 (CAPA) — fundo #323E30, texto branco -----
    s = prs.slides.add_slide(layout_branco)
    _definir_fundo(s, VERDE_ESCURO)
    adicionar_caixa(
        s, MARGEM, Inches(1.6), LARGURA_CONTEUDO, Inches(4.2),
        [
            {
                "texto": "Sua fazenda inteira em um lugar so. Ou ainda "
                         "espalhada no caderno, na planilha e na memoria?",
                "tamanho": 38, "negrito": True,
            },
        ],
        cor=BRANCO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
    )
    adicionar_caixa(
        s, MARGEM, Inches(5.9), LARGURA_CONTEUDO, Inches(0.9),
        [{"texto": "Conheca o GestSilo", "tamanho": 24, "cor": VERDE_CAMPO}],
        cor=VERDE_CAMPO, align=PP_ALIGN.CENTER,
    )
    adicionar_placeholder_foto(s, BRANCO)
    adicionar_logo(s, BRANCO)
    fundos.append(("Slide 1 (Capa)", "#323E30 (verde escuro) / texto branco"))

    # ----- SLIDE 2 — fundo branco, texto #323E30 -----
    s = prs.slides.add_slide(layout_branco)
    _definir_fundo(s, BRANCO)
    adicionar_caixa(
        s, MARGEM, MARGEM, LARGURA_CONTEUDO, Inches(1.4),
        [{"texto": "Voce conhece essas cenas", "tamanho": 32, "negrito": True}],
        cor=VERDE_ESCURO,
    )
    adicionar_caixa(
        s, MARGEM, Inches(2.8), LARGURA_CONTEUDO, Inches(5.5),
        [
            {"texto": "A silagem que acaba antes do previsto.", "tamanho": 20},
            {"texto": "O custo da lavoura que ninguem soube ao certo.",
             "tamanho": 20},
            {"texto": "O peso anotado num caderno que sumiu.", "tamanho": 20},
            {"texto": "E no fim do ano, a pergunta sem resposta: deu lucro "
                      "ou nao?", "tamanho": 20},
        ],
        cor=VERDE_ESCURO, espacamento_paragrafos=Pt(18),
    )
    adicionar_logo(s, VERDE_ESCURO)
    fundos.append(("Slide 2", "#FFFFFF (branco) / texto #323E30"))

    # ----- SLIDE 3 — fundo #323E30, texto branco -----
    s = prs.slides.add_slide(layout_branco)
    _definir_fundo(s, VERDE_ESCURO)
    adicionar_caixa(
        s, MARGEM, MARGEM, LARGURA_CONTEUDO, Inches(2.0),
        [{"texto": "O problema nao e falta de trabalho. E falta de controle.",
          "tamanho": 32, "negrito": True}],
        cor=BRANCO,
    )
    adicionar_caixa(
        s, MARGEM, Inches(3.4), LARGURA_CONTEUDO, Inches(4.5),
        [
            {"texto": "A informacao existe — mas vive espalhada na memoria, "
                      "em papeis e em planilhas que ninguem cruza.",
             "tamanho": 20},
            {"texto": "E o que nao se mede, se perde: em dinheiro, em tempo e "
                      "em tranquilidade.", "tamanho": 20},
        ],
        cor=BRANCO, espacamento_paragrafos=Pt(18),
    )
    adicionar_logo(s, BRANCO)
    fundos.append(("Slide 3", "#323E30 (verde escuro) / texto branco"))

    # ----- SLIDE 4 — fundo branco, texto #323E30 -----
    s = prs.slides.add_slide(layout_branco)
    _definir_fundo(s, BRANCO)
    adicionar_caixa(
        s, MARGEM, MARGEM, LARGURA_CONTEUDO, Inches(1.4),
        [{"texto": "Tudo o que acontece dentro da porteira", "tamanho": 32,
          "negrito": True}],
        cor=VERDE_ESCURO,
    )
    adicionar_caixa(
        s, MARGEM, Inches(2.8), LARGURA_CONTEUDO, Inches(5.5),
        [
            {"texto": "Silos de silagem e estoque de volumoso", "tamanho": 20},
            {"texto": "Rebanho — leite e corte", "tamanho": 20},
            {"texto": "Lavouras, pastagens e custo real de cada area",
             "tamanho": 20},
            {"texto": "Financeiro, insumos, frota e mao de obra", "tamanho": 20},
            {"texto": "Balanco forrageiro e planejamento de silagem",
             "tamanho": 20},
        ],
        cor=VERDE_ESCURO, espacamento_paragrafos=Pt(16),
    )
    adicionar_logo(s, VERDE_ESCURO)
    fundos.append(("Slide 4", "#FFFFFF (branco) / texto #323E30"))

    # ----- SLIDE 5 — fundo #323E30, texto branco -----
    s = prs.slides.add_slide(layout_branco)
    _definir_fundo(s, VERDE_ESCURO)
    adicionar_caixa(
        s, MARGEM, MARGEM, LARGURA_CONTEUDO, Inches(1.4),
        [{"texto": "Feito para o campo de verdade", "tamanho": 32,
          "negrito": True}],
        cor=BRANCO,
    )
    adicionar_caixa(
        s, MARGEM, Inches(2.8), LARGURA_CONTEUDO, Inches(5.0),
        [
            {"texto": "Funciona mesmo onde a internet e fraca: voce registra "
                      "a retirada no silo sem sinal, e o sistema sincroniza "
                      "depois.", "tamanho": 20},
            {"texto": "Tela legivel ao sol, operacao simples ate com luva.",
             "tamanho": 20},
        ],
        cor=BRANCO, espacamento_paragrafos=Pt(18),
    )
    adicionar_logo(s, BRANCO)
    fundos.append(("Slide 5", "#323E30 (verde escuro) / texto branco"))

    # ----- SLIDE 6 (CTA) — fundo #738D45, texto branco -----
    s = prs.slides.add_slide(layout_branco)
    _definir_fundo(s, VERDE_CAMPO)
    adicionar_caixa(
        s, MARGEM, Inches(1.4), LARGURA_CONTEUDO, Inches(1.6),
        [{"texto": "Comece de graca.", "tamanho": 44, "negrito": True}],
        cor=BRANCO, align=PP_ALIGN.CENTER,
    )
    adicionar_caixa(
        s, MARGEM, Inches(3.2), LARGURA_CONTEUDO, Inches(1.8),
        [{"texto": "Plano gratuito, sem cartao de credito. O plano cresce "
                   "junto com a sua fazenda.", "tamanho": 22}],
        cor=BRANCO, align=PP_ALIGN.CENTER,
    )
    adicionar_caixa(
        s, MARGEM, Inches(5.1), LARGURA_CONTEUDO, Inches(0.9),
        [{"texto": "gestsilo.com.br", "tamanho": 26, "negrito": True}],
        cor=BRANCO, align=PP_ALIGN.CENTER,
    )
    adicionar_placeholder_foto(s, BRANCO)
    adicionar_logo(s, BRANCO)
    fundos.append(("Slide 6 (CTA)", "#738D45 (verde campo) / texto branco"))

    return prs, fundos


def main():
    prs, fundos = construir()
    nome_arquivo = "carrossel_1_apresentacao.pptx"
    prs.save(nome_arquivo)

    print(f"Arquivo gerado: {nome_arquivo}")
    print(f"Total de slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
    print("-" * 60)
    for titulo, fundo in fundos:
        print(f"  {titulo:18s} -> fundo {fundo}")


if __name__ == "__main__":
    main()
